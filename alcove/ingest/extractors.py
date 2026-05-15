from __future__ import annotations

import csv
import io
import json
from pathlib import Path
from typing import List


def extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [p.extract_text() or "" for p in reader.pages]
    return "\n".join(pages)


def extract_epub(path: Path) -> str:
    from ebooklib import epub
    from bs4 import BeautifulSoup

    book = epub.read_epub(str(path))
    texts: List[str] = []
    for item in book.get_items():
        if item.get_type() == 9:  # DOCUMENT
            soup = BeautifulSoup(item.get_body_content(), "html.parser")
            texts.append(soup.get_text(" "))
    return "\n".join(texts)


def extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    return soup.get_text(separator=" ", strip=True)


def extract_md(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_rst(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_csv(path: Path, delimiter: str = ",") -> str:
    text = path.read_text(encoding="utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    return " ".join(cell for row in reader for cell in row)


def extract_tsv(path: Path) -> str:
    return extract_csv(path, delimiter="\t")


def extract_json(path: Path) -> str:
    data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(data, ensure_ascii=False)


def extract_jsonl(path: Path) -> str:
    lines = []
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        line = line.strip()
        if line:
            lines.append(json.dumps(json.loads(line), ensure_ascii=False))
    return "\n".join(lines)


def extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as e:
        raise ImportError("python-docx is required for .docx support: pip install python-docx") from e

    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError("python-pptx is required for .pptx support: pip install python-pptx") from e

    presentation = Presentation(str(path))
    texts: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
            if text:
                texts.append(text)
    return "\n".join(texts)


def extract_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as e:
        raise ImportError("striprtf is required for .rtf support: pip install 'alcove-search[rtf]'") from e

    return rtf_to_text(path.read_text(encoding="utf-8", errors="ignore"))
