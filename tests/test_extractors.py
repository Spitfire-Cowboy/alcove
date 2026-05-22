import json
from pathlib import Path

import pytest

from alcove.ingest.extractors import extract_csv, extract_json, extract_jsonl, extract_tsv


def test_extract_csv_joins_cells(tmp_path):
    f = tmp_path / "data.csv"
    f.write_text("a,b\n1,2\n")
    assert extract_csv(f) == "a b 1 2"


def test_extract_json_roundtrip(tmp_path):
    f = tmp_path / "data.json"
    f.write_text('{"k": "v", "n": 1}')
    out = extract_json(f)
    assert '"k": "v"' in out and '"n": 1' in out


def test_extract_jsonl_preserves_lines(tmp_path):
    f = tmp_path / "data.jsonl"
    f.write_text('{"a":1}\n{"b":2}\n')
    out = extract_jsonl(f)
    assert '{"a": 1}' in out and '{"b": 2}' in out


def test_extract_docx_returns_paragraph_text(tmp_path):
    docx = pytest.importorskip("docx", reason="python-docx not installed")
    from alcove.ingest.extractors import extract_docx

    doc = docx.Document()
    doc.add_paragraph("Alcove document extraction test.")
    f = tmp_path / "sample.docx"
    doc.save(str(f))

    result = extract_docx(f)
    assert "Alcove document extraction test." in result


def test_extract_odt_returns_paragraph_text(tmp_path):
    pytest.importorskip("odf.opendocument", reason="odfpy not installed")
    from odf import teletype
    from odf.opendocument import OpenDocumentText
    from odf.text import H, P

    from alcove.ingest.extractors import extract_odt

    doc = OpenDocumentText()
    heading = H(outlinelevel=1)
    teletype.addTextToElement(heading, "Alcove ODT Heading")
    doc.text.addElement(heading)
    paragraph = P()
    teletype.addTextToElement(paragraph, "ODT extraction body text.")
    doc.text.addElement(paragraph)

    f = tmp_path / "sample.odt"
    doc.save(str(f))

    result = extract_odt(f)
    assert "Alcove ODT Heading" in result
    assert "ODT extraction body text." in result


def _write_sample_pptx(path: Path, lines: list[str]) -> None:
    pptx = pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    frame = box.text_frame
    frame.text = lines[0]
    for line in lines[1:]:
        frame.add_paragraph().text = line
    presentation.save(str(path))


def test_extract_pptx_returns_slide_text(tmp_path):
    from alcove.ingest.extractors import extract_pptx

    f = tmp_path / "sample.pptx"
    _write_sample_pptx(f, ["PPT heading", "Bullet text"])

    result = extract_pptx(f)
    assert "PPT heading" in result
    assert "Bullet text" in result


def test_extract_pptx_raises_helpful_import_error(tmp_path, monkeypatch):
    from alcove.ingest import extractors

    f = tmp_path / "sample.pptx"
    f.write_bytes(b"not-a-real-pptx")

    real_import = __import__

    def blocked_import(name, *args, **kwargs):
        if name == "pptx":
            raise ImportError("missing")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", blocked_import)
    with pytest.raises(ImportError, match="python-pptx is required"):
        extractors.extract_pptx(f)


def test_extract_rtf_returns_plain_text(tmp_path):
    pytest.importorskip("striprtf.striprtf", reason="striprtf not installed")
    from alcove.ingest.extractors import extract_rtf

    f = tmp_path / "sample.rtf"
    f.write_text(r"{\rtf1\ansi Alcove {\b RTF} extractor test.\par Second paragraph.}")

    result = extract_rtf(f)
    assert "Alcove" in result
    assert "RTF" in result
    assert "Second paragraph." in result
    assert "\\b" not in result


def test_extract_rtf_requires_striprtf(tmp_path, monkeypatch):
    from alcove.ingest import extractors

    f = tmp_path / "sample.rtf"
    f.write_text(r"{\rtf1\ansi missing dependency test}")

    real_import = __import__

    def blocked_import(name, *args, **kwargs):
        if name in {"striprtf", "striprtf.striprtf"}:
            raise ImportError("missing")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", blocked_import)
    with pytest.raises(ImportError, match=r"striprtf is required.*alcove-search\[rtf\]"):
        extractors.extract_rtf(f)


def test_extract_xlsx_returns_cell_text(tmp_path):
    openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")
    from alcove.ingest.extractors import extract_xlsx

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Name"
    ws["B1"] = "Score"
    ws["A2"] = "Alice"
    ws["B2"] = 42
    f = tmp_path / "sample.xlsx"
    wb.save(str(f))

    result = extract_xlsx(f)
    assert "Name" in result
    assert "Alice" in result
    assert "42" in result
    assert "Score" in result


def test_extract_xlsx_multiple_sheets(tmp_path):
    openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")
    from alcove.ingest.extractors import extract_xlsx

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1["A1"] = "hello"
    ws2 = wb.create_sheet("Sheet2")
    ws2["A1"] = "world"
    f = tmp_path / "multi.xlsx"
    wb.save(str(f))

    result = extract_xlsx(f)
    assert "hello" in result
    assert "world" in result


def test_extract_xlsx_skips_none_cells(tmp_path):
    openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")
    from alcove.ingest.extractors import extract_xlsx

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "alpha"
    ws["B1"] = None
    ws["C1"] = "omega"
    f = tmp_path / "sparse.xlsx"
    wb.save(str(f))

    result = extract_xlsx(f)
    assert "alpha" in result
    assert "omega" in result


def test_extract_odt_requires_odfpy(tmp_path, monkeypatch):
    from alcove.ingest import extractors

    f = tmp_path / "sample.odt"
    f.write_bytes(b"not-a-real-odt")

    real_import = __import__

    def blocked_import(name, *args, **kwargs):
        if name in {"odf", "odf.opendocument", "odf.text"}:
            raise ImportError("missing")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", blocked_import)
    with pytest.raises(ImportError, match=r"odfpy is required.*alcove-search\[odt\]"):
        extractors.extract_odt(f)


def test_extract_xlsx_requires_openpyxl(tmp_path, monkeypatch):
    from alcove.ingest import extractors

    f = tmp_path / "sample.xlsx"
    f.write_bytes(b"not-a-real-xlsx")

    real_import = __import__

    def blocked_import(name, *args, **kwargs):
        if name == "openpyxl":
            raise ImportError("missing")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", blocked_import)
    with pytest.raises(ImportError, match=r"openpyxl is required.*alcove-search\[xlsx\]"):
        extractors.extract_xlsx(f)


def test_pipeline_dispatch_includes_html(tmp_path):
    """Pipeline routes .html files through the extractor without error."""
    from alcove.ingest.pipeline import run

    f = tmp_path / "doc.html"
    f.write_text("<p>Semantic search content</p>")
    out = tmp_path / "chunks.jsonl"
    n = run(raw_dir=str(tmp_path), out_file=str(out))
    assert n >= 1
    records = [json.loads(line) for line in out.read_text().splitlines()]
    assert any("Semantic search content" in r["chunk"] for r in records)


def test_pipeline_dispatch_includes_md(tmp_path):
    """Pipeline routes .md files through the extractor without error."""
    from alcove.ingest.pipeline import run

    f = tmp_path / "notes.md"
    f.write_text("# Notes\n\nLocal retrieval is fast.")
    out = tmp_path / "chunks.jsonl"
    n = run(raw_dir=str(tmp_path), out_file=str(out))
    assert n >= 1
    records = [json.loads(line) for line in out.read_text().splitlines()]
    assert any("Local retrieval is fast." in r["chunk"] for r in records)


def test_pipeline_dispatch_includes_pptx(tmp_path):
    """Pipeline routes .pptx files through the extractor without error."""
    from alcove.ingest.pipeline import run

    f = tmp_path / "deck.pptx"
    _write_sample_pptx(f, ["Deck heading", "Deck body text"])
    out = tmp_path / "chunks.jsonl"
    n = run(raw_dir=str(tmp_path), out_file=str(out))
    assert n >= 1
    records = [json.loads(line) for line in out.read_text().splitlines()]
    assert any("Deck heading" in r["chunk"] for r in records)


def test_pipeline_dispatch_includes_odt(tmp_path):
    pytest.importorskip("odf.opendocument", reason="odfpy not installed")
    from odf import teletype
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    from alcove.ingest.pipeline import run

    doc = OpenDocumentText()
    paragraph = P()
    teletype.addTextToElement(paragraph, "Pipeline ODT content.")
    doc.text.addElement(paragraph)

    f = tmp_path / "sample.odt"
    doc.save(str(f))
    out = tmp_path / "chunks.jsonl"

    n = run(raw_dir=str(tmp_path), out_file=str(out))

    assert n >= 1
    records = [json.loads(line) for line in out.read_text().splitlines()]
    assert any("Pipeline ODT content." in r["chunk"] for r in records)


def test_pipeline_dispatch_includes_xlsx(tmp_path):
    openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")
    from alcove.ingest.pipeline import run

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "spreadsheet content for search"
    f = tmp_path / "sheet.xlsx"
    wb.save(str(f))
    out = tmp_path / "chunks.jsonl"

    n = run(raw_dir=str(tmp_path), out_file=str(out))

    assert n >= 1
    records = [json.loads(line) for line in out.read_text().splitlines()]
    assert any("spreadsheet content for search" in r["chunk"] for r in records)



def test_pipeline_dispatch_includes_rtf(tmp_path):
    pytest.importorskip("striprtf.striprtf", reason="striprtf not installed")
    from alcove.ingest.pipeline import run

    f = tmp_path / "sample.rtf"
    f.write_text(r"{\rtf1\ansi Pipeline {\b RTF} content.}")
    out = tmp_path / "chunks.jsonl"
    n = run(raw_dir=str(tmp_path), out_file=str(out))

    assert n >= 1
    records = [json.loads(line) for line in out.read_text().splitlines()]
    assert any("Pipeline RTF content." in r["chunk"] for r in records)

def test_extract_tsv_tab_separated(tmp_path):
    """extract_tsv correctly parses tab-delimited files."""
    f = tmp_path / "data.tsv"
    f.write_text("x\ty\n3\t4\n")
    assert extract_tsv(f) == "x y 3 4"
